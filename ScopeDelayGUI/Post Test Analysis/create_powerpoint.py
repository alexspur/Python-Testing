"""
Create PowerPoint Presentation from Marx Generator Test Figures
"""

from pptx import Presentation
from pptx.util import Inches, Pt
# from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN
import os
import re
from datetime import datetime
from PIL import Image

#%% ===================== CONFIGURATION =====================
# Set the date for the presentation filename
test_date = "12.11.25"

# Base folder containing the Figures folder
base_folder = r"C:\Users\ESpurbeck\Desktop\LANL Project\Marx Testing\Testing\12.11.25"

# Figures folder
figures_folder = os.path.join(base_folder, "Figures")

# Output PowerPoint filename
output_filename = f"Marx_Test_Results_{test_date}.pptx"
output_path = os.path.join(base_folder, output_filename)

# Presentation title
presentation_title = f"Marx Generator Test Results - {test_date}"

#%% ===================== CREATE PRESENTATION =====================
print(f"Creating PowerPoint presentation...")
print(f"Figures folder: {figures_folder}")

# Create presentation with widescreen dimensions (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Get all PNG files in the figures folder
png_files = [f for f in os.listdir(figures_folder) if f.endswith('.png')]

# Sort files by test number
def extract_test_num(filename):
    match = re.search(r'Test_(\d+)', filename)
    return int(match.group(1)) if match else 999

png_files.sort(key=extract_test_num)

print(f"Found {len(png_files)} figures")

#%% ===================== ADD TITLE SLIDE =====================
title_slide_layout = prs.slide_layouts[6]  # Blank slide
slide = prs.slides.add_slide(title_slide_layout)

# Add title text box
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = presentation_title
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.name = "Times New Roman"
title_para.alignment = PP_ALIGN.CENTER

# Add subtitle with date
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
subtitle_para.font.size = Pt(24)
subtitle_para.font.name = "Times New Roman"
subtitle_para.alignment = PP_ALIGN.CENTER

#%% ===================== ADD FIGURE SLIDES =====================
#%% ===================== ADD FIGURE SLIDES =====================
blank_slide_layout = prs.slide_layouts[6]  # Blank slide

# Fixed image dimensions
img_width = Inches(9.23)
img_height = Inches(6.15)

for i, png_file in enumerate(png_files):
    print(f"Adding slide {i+2}: {png_file}")
    
    # Add slide
    slide = prs.slides.add_slide(blank_slide_layout)
    
    img_path = os.path.join(figures_folder, png_file)
    
    # Center the image on the slide
    left = (prs.slide_width - img_width) / 2
    top = (prs.slide_height - img_height) / 2
    
    slide.shapes.add_picture(img_path, left, top, width=img_width, height=img_height)

#%% ===================== SAVE PRESENTATION =====================
prs.save(output_path)
print(f"\n========================================")
print(f"Presentation saved to:")
print(f"{output_path}")
print(f"========================================")
print(f"Total slides: {len(png_files) + 1} (1 title + {len(png_files)} figures)")
